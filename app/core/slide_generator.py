import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime

class SlideGenerator:
    """Generates professional PowerPoint presentations from analysis results."""
    
    def __init__(self, brand_color=(124, 58, 237)): # primary-600 #7c3aed
        self.brand_color = RGBColor(*brand_color)
        self.text_color = RGBColor(30, 41, 59) # slate-800
        self.muted_color = RGBColor(100, 116, 139) # slate-500

    def _add_styled_slide(self, prs, title_text):
        """Adds a slide with a consistent header and branding."""
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Style Title
        title = slide.shapes.title
        title.text = title_text
        title_text_frame = title.text_frame
        p = title_text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.color.rgb = self.brand_color
        
        # Add a subtle footer line
        shapes = slide.shapes
        line = shapes.add_shape(
            1, # Rectangle for a line
            Inches(0.5), Inches(1.1), Inches(9), Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.brand_color
        line.line.visible = False
        
        return slide

    def generate(self, data, filename):
        """Creates a PPTX presentation from analysis data."""
        prs = Presentation()
        
        # 1. Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.shapes.placeholders[1]
        
        title.text = "Executive Strategic Analysis"
        subtitle.text = f"Report: {filename}\nGenerated on: {datetime.date.today().strftime('%B %d, %Y')}\nAI Market Analyst Pro"
        
        # Style Title Slide
        title.text_frame.paragraphs[0].font.color.rgb = self.brand_color
        title.text_frame.paragraphs[0].font.bold = True
        
        # 2. Executive Summary (Overview)
        summary_slide = self._add_styled_slide(prs, "Executive Summary")
        tf = summary_slide.placeholders[1].text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get('current_business_situation', 'No overview provided.')
        p.font.size = Pt(18)
        p.space_after = Pt(10)

        # 3. Strategic Strengths
        strengths_slide = self._add_styled_slide(prs, "Strategic Strengths")
        tf = strengths_slide.placeholders[1].text_frame
        tf.word_wrap = True
        strong_points = data.get('strong_points', [])
        
        # Adjust density for long lists
        font_size = 16 if len(strong_points) <= 6 else (14 if len(strong_points) <= 8 else 12)
        
        for strength in strong_points:
            p = tf.add_paragraph()
            # Handle both string and dict formats
            p.text = strength.get('content', str(strength)) if isinstance(strength, dict) else str(strength)
            p.level = 0
            p.space_after = Pt(6)
            p.font.size = Pt(font_size)

        # 4. Critical Improvements
        weakness_slide = self._add_styled_slide(prs, "Areas for Improvement")
        tf = weakness_slide.placeholders[1].text_frame
        tf.word_wrap = True
        weak_points = data.get('weak_points', [])
        
        font_size = 16 if len(weak_points) <= 6 else (14 if len(weak_points) <= 8 else 12)
        
        for point in weak_points:
            p = tf.add_paragraph()
            p.text = point.get('content', str(point)) if isinstance(point, dict) else str(point)
            p.level = 0
            p.space_after = Pt(6)
            p.font.size = Pt(font_size)

        # 5. Strategic Roadmap
        roadmap_slide = self._add_styled_slide(prs, "Strategic Roadmap")
        tf = roadmap_slide.placeholders[1].text_frame
        tf.word_wrap = True
        moves = data.get('next_strategic_moves', [])
        
        font_size = 18 if len(moves) <= 5 else (16 if len(moves) <= 7 else 14)
        
        for move in moves:
            p = tf.add_paragraph()
            action = move.get('action', str(move)) if isinstance(move, dict) else str(move)
            priority = f" [{move.get('priority', '')}]" if isinstance(move, dict) and move.get('priority') else ""
            p.text = f"{action}{priority}"
            p.level = 0
            p.font.bold = True
            p.space_before = Pt(8)
            p.font.size = Pt(font_size)
        
        # New: Strategic Shifts (Comparison only)
        if data.get('is_comparison') and data.get('comparison_delta'):
            shifts_slide = self._add_styled_slide(prs, "Strategic Shifts (Report Comparison)")
            tf = shifts_slide.placeholders[1].text_frame
            for shift in data.get('comparison_delta', []):
                p = tf.add_paragraph()
                p.text = shift
                p.level = 0
                p.space_after = Pt(8)
                p.font.size = Pt(16)
                p.font.bold = True

        # 6. Live Market Insights (Optional)
        if data.get('market_context'):
            market_slide = self._add_styled_slide(prs, "Live Market Intelligence")
            tf = market_slide.placeholders[1].text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data['market_context']
            p.font.size = Pt(16)
            p.font.italic = True

        # Save to BytesIO
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io
